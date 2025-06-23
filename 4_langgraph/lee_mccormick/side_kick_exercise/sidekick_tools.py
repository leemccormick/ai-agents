from playwright.async_api import async_playwright
from langchain_community.agent_toolkits import PlayWrightBrowserToolkit
from dotenv import load_dotenv
import os
import requests
from langchain.agents import Tool
from langchain_community.agent_toolkits import FileManagementToolkit
from langchain_community.tools.wikipedia.tool import WikipediaQueryRun
from langchain_experimental.tools import PythonREPLTool
from langchain_community.utilities import GoogleSerperAPIWrapper
from langchain_community.utilities.wikipedia import WikipediaAPIWrapper
from sendgrid.helpers.mail import Mail, Email, To, Content, Attachment, FileContent, FileName, FileType, Disposition
from pptx import Presentation
import base64
from sidekick_models import PresenterOutput
import base64
from typing import Optional
from sendgrid import SendGridAPIClient
from sendgrid.helpers.mail import (
    Mail, Email, To, Content, Attachment,
    FileContent, FileName, FileType, Disposition
)

load_dotenv(override=True)
pushover_token = os.getenv("PUSHOVER_TOKEN")
pushover_user = os.getenv("PUSHOVER_USER")
pushover_url = "https://api.pushover.net/1/messages.json"
serper = GoogleSerperAPIWrapper()
sendgrid_api_key = os.getenv("SENDGRID_API_KEY")

async def playwright_tools():
    playwright = await async_playwright().start()
    browser = await playwright.chromium.launch(headless=False)
    toolkit = PlayWrightBrowserToolkit.from_browser(async_browser=browser)
    return toolkit.get_tools(), browser, playwright

def push(text: str):
    """Send a push notification to the user"""
    requests.post(pushover_url, data = {"token": pushover_token, "user": pushover_user, "message": text})
    return "success"

def get_file_tools():
    toolkit = FileManagementToolkit(root_dir="ai-info")
    return toolkit.get_tools()

def send_email(body: str, pptx_path: Optional[str] = None):
    """Send an email with the given body and optionally attach a PowerPoint file."""

    from_email = Email("leemccormick.developer@gmail.com")
    to_email = To("lovebirdsmccormick@gmail.com")
    subject = "Sales email" + (" with presentation" if pptx_path else "")
    content = Content("text/plain", body)

    # Compose email
    mail = Mail(from_email, to_email, subject, content)

    # Attach PowerPoint if provided
    if pptx_path:
        with open(pptx_path, "rb") as f:
            file_data = f.read()
        encoded_file = base64.b64encode(file_data).decode()

        attachment = Attachment(
            FileContent(encoded_file),
            FileName(os.path.basename(pptx_path)),
            FileType("application/vnd.openxmlformats-officedocument.presentationml.presentation"),
            Disposition("attachment")
        )
        mail.attachment = attachment

    # Send email
    sg = SendGridAPIClient(api_key=sendgrid_api_key)
    response = sg.client.mail.send.post(request_body=mail.get())

    return {
        "status": response.status_code,
        "message": response.body.decode() if hasattr(response.body, "decode") else str(response.body)
    }


def presentation(data: PresenterOutput):
    # Create presentation object
    prs = Presentation()
    
    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]  # Layout 0 = Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = data.slide1Title
    subtitle.text = data.slide2Subtitle
    
    # Slide 2: Content Slide with Bullet Points
    bullet_slide_layout = prs.slide_layouts[1]  # Layout 1 = Title and Content
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = data.slide2Title
    content.text = data.slide2Content
    
    # Slide 3: Two Column Layout
    two_column_layout = prs.slide_layouts[3]  # Layout 3 = Two Content
    slide = prs.slides.add_slide(two_column_layout)
    title = slide.shapes.title
    left_content = slide.placeholders[1]
    right_content = slide.placeholders[2]
    
    title.text = data.slide3Title
    left_content.text = data.slide3LeftContent
    right_content.text = data.slide3RightContent
    
    # Slide 4: Thank You Slide
    thank_you_layout = prs.slide_layouts[0]  # Use title layout again
    slide = prs.slides.add_slide(thank_you_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = data.slide4Title
    subtitle.text = data.slide4Subtitle

    # Define the folder and file path
    folder_name = "ai-presentation"
    file_name = "{data.file_name}.pptx"
    output_path = os.path.join(folder_name, file_name)

    # Ensure the folder exists
    os.makedirs(folder_name, exist_ok=True)

    # Save the presentation to the desired path
    prs.save(output_path)

    print(f"Presentation saved to: {output_path}")
    return output_path


async def search_tools():
    serper_search =Tool(
        name="search",
        func=serper.run,
        description="Use this tool when you want to get the results of an online web search"
    )

    wikipedia = WikipediaAPIWrapper()
    wiki_tool = WikipediaQueryRun(api_wrapper=wikipedia)
    python_repl = PythonREPLTool()

    return [serper_search, python_repl,  wiki_tool]

async def write_tools():
    file_tools = get_file_tools()
    presentation_tool = Tool(
        name="presentation_tool",
        func=presentation,
        description="Use this tool to create presentation as .pptx file"
    )

    return file_tools + [presentation_tool]

async def report_tools():
    push_tool = Tool(
        name="send_push_notification", 
        func=push, 
        description="Use this tool when you want to send a push notification"
    )

    send_email_tool = Tool(
        name="send_email",
        func=send_email,
        description="Use this tool when you want to send an email"
    )

    return [push_tool, send_email_tool]
